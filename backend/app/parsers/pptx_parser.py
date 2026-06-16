from __future__ import annotations

from pathlib import Path
from typing import List, Dict

try:
    from pptx import Presentation
except Exception as exc:  # pragma: no cover - runtime guard
    Presentation = None


def parse_pptx(path: Path) -> List[Dict]:
    """
    Extract slide-level content from a .pptx file.
    Returns a list of dicts: {title, bullets, notes, full_text}
    """
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is not available. Install it with: pip install python-pptx"
        )

    prs = Presentation(str(path))
    slides = []

    for slide in prs.slides:
        title = ""
        bullets: List[str] = []
        body_lines: List[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            # Heuristic: if the shape is the title placeholder, treat as title.
            if shape.is_placeholder:
                from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore

                placeholder_type = shape.placeholder_format.type
                if placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    title = text
                    continue

            # Otherwise collect paragraphs as bullets / body text.
            for paragraph in shape.text_frame.paragraphs:
                line = paragraph.text.strip()
                if line:
                    body_lines.append(line)
                    # Treat top-level paragraphs with bullets as bullets.
                    if paragraph.level == 0:
                        bullets.append(line)

        # Fallback: if no title placeholder was found, treat the first body line as title.
        if not title and body_lines:
            title = body_lines[0]
            body_lines = body_lines[1:]
            if bullets and bullets[0] == title:
                bullets = bullets[1:]

        # If no explicit bullets were captured, use the remaining body lines.
        if not bullets and body_lines:
            bullets = body_lines

        # Speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        full_text = "\n".join([title] + bullets)

        slides.append(
            {
                "title": title,
                "bullets": bullets,
                "notes": notes,
                "full_text": full_text,
            }
        )

    return slides
