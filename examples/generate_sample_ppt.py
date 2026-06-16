"""
Generate a sample pitch deck for the Pharos AI Agent Economy demo.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


SLIDES = [
    {
        "title": "Pharos Agent Economy",
        "bullets": [
            "High-performance EVM Layer 1 for Web2/Web3 fusion",
            "Built for RealFi, RWA, stablecoins and cross-chain interoperability",
            "Designed to support millions of autonomous AI Agents",
        ],
    },
    {
        "title": "The Problem: Agents Cannot Explain Themselves",
        "bullets": [
            "Wallet analyzers, swaps and monitors only return JSON or text",
            "DAO members and users struggle to read raw on-chain data",
            "Project teams lack reusable, verifiable presentation tools",
        ],
    },
    {
        "title": "PitchFlow: Agent Expression Layer",
        "bullets": [
            "Reusable Skill: any Agent can call Deck-to-Video API",
            "Generates narrated MP4, subtitles, cover and metadata",
            "Content hash registered on Pharos for provenance",
        ],
    },
    {
        "title": "Composable On-chain Economy",
        "bullets": [
            "Wallet Analyzer → Report PPT → Narrated Video",
            "Payment / tipping via PaymentGate smart contract",
            "Distribution to Telegram, Discord, X and DAO forums",
        ],
    },
    {
        "title": "Roadmap & Call to Action",
        "bullets": [
            "Skill stage: standard API + Pharos registry integration",
            "Agent Arena: full PitchFlow Agent with payments and social distribution",
            "Join us in building the communication infrastructure of Agent Economy",
        ],
    },
]


def add_slide(prs, title, bullets):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True

    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(4.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(26)
        p.space_after = Pt(14)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        add_slide(prs, slide_data["title"], slide_data["bullets"])

    out = Path(__file__).with_name("pharos-agent-weekly-report.pptx")
    prs.save(out)
    print(f"Created sample deck: {out}")


if __name__ == "__main__":
    main()
